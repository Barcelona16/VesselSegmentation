from pptx import Presentation
 
# 创建幻灯片 ------
prs = Presentation()
title_slide_layout = prs.slide_layouts[0] #2 line
information_layout = prs.slide_layouts[1] # n lines for information

slide = prs.slides.add_slide(title_slide_layout) 
title = slide.shapes.title
subtitle = slide.placeholders[1]
# 设置标题和副标题
title.text = "本周工作内容"
subtitle.text = "Deavan"
slide = prs.slides.add_slide(information_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = ".tre文件"
subtitle.text = "1.对tre文件读取代码进行封装（draw.py）\n2.学习MicroTumordetection-NeuralNet-PixelClassifier"

slide = prs.slides.add_slide(information_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "对于血管数据文件"
subtitle.text = "1.学习simpleITK\n2.https://pyscience.wordpress.com/2014/11/02/multi-modal-image-segmentation-with-python-simpleitk \n3.https://pyscience.wordpress.com/2014/10/19/image-segmentation-with-python-and-simpleitk/"

slide = prs.slides.add_slide(title_slide_layout) 
title = slide.shapes.title
subtitle = slide.placeholders[1]
# 设置标题和副标题
title.text = "问题"
subtitle.text = "1.tre文件读取方法的不确定\n2.方法调研"

slide = prs.slides.add_slide(title_slide_layout) 
title = slide.shapes.title
subtitle = slide.placeholders[1]
# 设置标题和副标题
title.text = "下周计划"
subtitle.text = "1.继续本周未完成的论文方法学习\n2.数据预处理\n3.方法学习"

 
prs.save("week2.pptx")
